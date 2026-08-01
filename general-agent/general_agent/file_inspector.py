"""Deterministic bounded previews for common uploaded file formats."""

from __future__ import annotations

import json
from itertools import islice
from pathlib import Path
from typing import Any

import pandas as pd
import yaml
from charset_normalizer import from_bytes
from docx import Document
from pypdf import PdfReader
from pptx import Presentation

from general_agent.config import Settings

_TEXT_EXTENSIONS = {
    ".txt", ".md", ".rst", ".py", ".js", ".ts", ".tsx", ".jsx",
    ".java", ".c", ".cc", ".cpp", ".h", ".hpp", ".go", ".rs",
    ".rb", ".php", ".sh", ".zsh", ".fish", ".sql", ".toml",
    ".ini", ".cfg", ".xml", ".html", ".css", ".log",
}


def inspect_path(path: Path, settings: Settings) -> dict[str, Any]:
    suffix = path.suffix.lower()
    common = {"name": path.name, "size_bytes": path.stat().st_size, "format": suffix.lstrip(".") or "unknown"}
    if suffix == ".pdf":
        return {**common, **_inspect_pdf(path, settings)}
    if suffix == ".docx":
        return {**common, **_inspect_docx(path, settings)}
    if suffix == ".pptx":
        return {**common, **_inspect_pptx(path, settings)}
    if suffix in {".xls", ".xlsx"}:
        return {**common, **_inspect_excel(path, settings)}
    if suffix in {".csv", ".tsv"}:
        return {**common, **_inspect_delimited(path, settings)}
    if suffix in {".json", ".yaml", ".yml"}:
        return {**common, **_inspect_structured(path, settings)}
    if suffix in _TEXT_EXTENSIONS or not suffix:
        return {**common, **_inspect_text(path, settings)}
    return {
        **common,
        "supported": False,
        "message": "This binary format can be stored and downloaded but has no deterministic preview.",
    }


def _bounded(text: str, limit: int) -> tuple[str, bool]:
    if len(text) <= limit:
        return text, False
    return text[:limit], True


def _inspect_pdf(path: Path, settings: Settings) -> dict[str, Any]:
    reader = PdfReader(str(path))
    previews = []
    for index, page in enumerate(reader.pages[: settings.max_inspect_pages], start=1):
        text, truncated = _bounded(page.extract_text() or "", settings.max_inspect_chars)
        previews.append({"page": index, "text": text, "truncated": truncated})
    scanned = not any(item["text"].strip() for item in previews)
    return {
        "supported": True,
        "page_count": len(reader.pages),
        "pages": previews,
        "preview_truncated": len(reader.pages) > settings.max_inspect_pages,
        "scanned_or_image_only": scanned,
        "message": "No embedded text was found; OCR is not available." if scanned else None,
    }


def _inspect_docx(path: Path, settings: Settings) -> dict[str, Any]:
    document = Document(str(path))
    paragraphs, truncated = _bounded(
        "\n".join(paragraph.text for paragraph in document.paragraphs),
        settings.max_inspect_chars,
    )
    tables = []
    for table in document.tables[: settings.max_inspect_sheets]:
        rows = [
            [cell.text for cell in row.cells[: settings.max_inspect_columns]]
            for row in table.rows[: settings.max_inspect_rows]
        ]
        tables.append(rows)
    return {
        "supported": True,
        "paragraph_text": paragraphs,
        "paragraph_text_truncated": truncated,
        "table_count": len(document.tables),
        "tables": tables,
    }


def _inspect_pptx(path: Path, settings: Settings) -> dict[str, Any]:
    presentation = Presentation(str(path))
    slides = []
    for index, slide in enumerate(
        islice(presentation.slides, settings.max_inspect_pages), start=1
    ):
        text_parts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text]
        notes = ""
        try:
            notes = slide.notes_slide.notes_text_frame.text
        except (AttributeError, ValueError):
            notes = ""
        text, truncated = _bounded("\n".join(text_parts), settings.max_inspect_chars)
        notes, notes_truncated = _bounded(notes, settings.max_inspect_chars)
        slides.append({"slide": index, "text": text, "notes": notes, "truncated": truncated or notes_truncated})
    return {
        "supported": True,
        "slide_count": len(presentation.slides),
        "slides": slides,
        "preview_truncated": len(presentation.slides) > settings.max_inspect_pages,
    }


def _frame_preview(frame: pd.DataFrame, settings: Settings) -> dict[str, Any]:
    bounded = frame.iloc[: settings.max_inspect_rows, : settings.max_inspect_columns]
    return {
        "rows": bounded.where(pd.notna(bounded), None).to_dict(orient="records"),
        "columns": [str(column) for column in bounded.columns],
        "row_count_previewed": len(bounded),
        "column_count_previewed": len(bounded.columns),
    }


def _inspect_excel(path: Path, settings: Settings) -> dict[str, Any]:
    workbook = pd.ExcelFile(path)
    sheets: dict[str, Any] = {}
    for sheet in workbook.sheet_names[: settings.max_inspect_sheets]:
        frame = workbook.parse(sheet, nrows=settings.max_inspect_rows)
        sheets[sheet] = _frame_preview(frame, settings)
    return {
        "supported": True,
        "sheet_names": workbook.sheet_names,
        "sheets": sheets,
        "preview_truncated": len(workbook.sheet_names) > settings.max_inspect_sheets,
    }


def _inspect_delimited(path: Path, settings: Settings) -> dict[str, Any]:
    separator = "\t" if path.suffix.lower() == ".tsv" else None
    frame = pd.read_csv(
        path,
        sep=separator,
        engine="python" if separator is None else "c",
        nrows=settings.max_inspect_rows,
    )
    return {"supported": True, **_frame_preview(frame, settings)}


def _decode_text(path: Path) -> tuple[str, str]:
    raw = path.read_bytes()
    match = from_bytes(raw).best()
    if match is None:
        raise ValueError("The file encoding could not be detected.")
    return str(match), match.encoding or "unknown"


def _inspect_text(path: Path, settings: Settings) -> dict[str, Any]:
    text, encoding = _decode_text(path)
    preview, truncated = _bounded(text, settings.max_inspect_chars)
    return {"supported": True, "encoding": encoding, "text": preview, "truncated": truncated}


def _inspect_structured(path: Path, settings: Settings) -> dict[str, Any]:
    text, encoding = _decode_text(path)
    if len(text) > settings.max_inspect_chars * 4:
        preview, truncated = _bounded(text, settings.max_inspect_chars)
        return {"supported": True, "encoding": encoding, "text": preview, "truncated": truncated, "parsed": False}
    value = json.loads(text) if path.suffix.lower() == ".json" else yaml.safe_load(text)
    rendered, truncated = _bounded(json.dumps(value, ensure_ascii=False, indent=2, default=str), settings.max_inspect_chars)
    return {"supported": True, "encoding": encoding, "text": rendered, "truncated": truncated, "parsed": True}
