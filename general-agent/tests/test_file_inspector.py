from __future__ import annotations

import json

import pandas as pd
from docx import Document
from pptx import Presentation
from reportlab.pdfgen import canvas

from general_agent.file_inspector import inspect_path


def test_text_structured_and_tabular_previews_are_bounded(settings) -> None:
    text = settings.workspace_root / "long.txt"
    text.write_text("x" * 1_000, encoding="utf-8")
    preview = inspect_path(text, settings)
    assert preview["supported"] is True
    assert preview["truncated"] is True
    assert len(preview["text"]) == settings.max_inspect_chars

    structured = settings.workspace_root / "data.json"
    structured.write_text(json.dumps({"hello": [1, 2]}), encoding="utf-8")
    assert inspect_path(structured, settings)["parsed"] is True

    csv = settings.workspace_root / "table.csv"
    pd.DataFrame({"a": range(8), "b": range(8)}).to_csv(csv, index=False)
    result = inspect_path(csv, settings)
    assert result["row_count_previewed"] == 3
    assert result["columns"] == ["a", "b"]


def test_office_and_pdf_previews(settings) -> None:
    docx = settings.workspace_root / "sample.docx"
    document = Document()
    document.add_paragraph("Word body")
    document.add_table(rows=1, cols=2).rows[0].cells[0].text = "Cell"
    document.save(docx)
    assert "Word body" in inspect_path(docx, settings)["paragraph_text"]

    pptx = settings.workspace_root / "slides.pptx"
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[1])
    slide.shapes.title.text = "Slide title"
    deck.save(pptx)
    assert inspect_path(pptx, settings)["slides"][0]["text"].startswith("Slide title")

    xlsx = settings.workspace_root / "book.xlsx"
    with pd.ExcelWriter(xlsx) as writer:
        pd.DataFrame({"value": [1, 2]}).to_excel(writer, sheet_name="Data", index=False)
    workbook = inspect_path(xlsx, settings)
    assert workbook["sheet_names"] == ["Data"]
    assert workbook["sheets"]["Data"]["rows"][0]["value"] == 1

    pdf = settings.workspace_root / "text.pdf"
    pdf_canvas = canvas.Canvas(str(pdf))
    pdf_canvas.drawString(72, 720, "Embedded text")
    pdf_canvas.save()
    assert "Embedded text" in inspect_path(pdf, settings)["pages"][0]["text"]

    scanned = settings.workspace_root / "scan.pdf"
    blank_canvas = canvas.Canvas(str(scanned))
    blank_canvas.showPage()
    blank_canvas.save()
    result = inspect_path(scanned, settings)
    assert result["scanned_or_image_only"] is True
    assert "OCR is not available" in result["message"]
